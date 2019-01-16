# -*- coding: utf-8; py-indent-offset:4 -*-
###############################################################################
#
# Copyright (C) Mineo Sudo
#
###############################################################################
from pptx import Presentation
import pandas as pd
import pd2ppt
import os

class pptx(object):
    '''
    pythonでパワポ作業を効率化するクラスです
    '''
    def __init__(self,file_path=""):
        self.file_path = file_path
        self.ppt = Presentation(file_path)
        self.foldr_path = os.path.dirname(self.file_path)
        self.file_name = os.path.basename(self.file_path)
        self.name, self.ext = os.path.splitext(self.file_name)
        self.df_ppt = self._make_df()
        self.csv_file_path = os.path.join(self.foldr_path,"temp",self.name+".csv" )

    def add_table_from_df(self,df):
        """pandasのデータフレームをパワポに流し込みます。先に流し込みたいpptxを指定しておくこと

        Arguments:
            df {pandasdataframe} -- パワポ化したいデータフレーム
        """

        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[0])
        # pd2ppt.df_to_table(slide,df, col_formatters=['', ',', '.2'],rounding=['', 3, ''])
        pd2ppt.df_to_table(slide,df)
        try:
            self.ppt.save(self.file_path)
            print("上書きしました")
        except:
            print("No file_path")

    def pptx_to_csv(self):
        """pptxのテキストとテーブルのデータをself.csv_file_pathに吐き出します
        """

        self.df_ppt.to_csv(self.csv_file_path ,index=False)
        print(self.csv_file_path+"にcsvを保存しました")

    def csv_to_pptx(self):
        """diffを見てcsvを編集した際に、これを実行すればpptxに反映される（テーブル未対応）
        """

        df_t = pd.read_csv(self.csv_file_path, dtype = 'object')
        li_henko = df_t.values.tolist()
        hata = 0
        for elm in li_henko:
            if elm[0].find("------slide ") == -1:
                taisho_shape = self._find_shape(int(elm[0]),int(elm[1]))
                if taisho_shape.has_text_frame:
                    if taisho_shape.text != elm[2] and isinstance(elm[2], str):
                        print("変更  :  " + taisho_shape.text +" → "+ str(elm[2]))
                        taisho_shape.text = elm[2]
                        hata = 1
            #     if taisho_shape.has_table:
            #             write_table(taisho_shape.table,elm)

        if hata == 0:
            print("変更なし")
        else:
            print("パワポを更新しました")
            self.ppt.save(self.file_path)

    def _make_df(self):
        """self.file_pathのパワポからpandas dataframeを作成します（文字と表のみ）
        
        Returns:
            pandas dataframe -- pptxから抽出されたdataframeです
        """
        i = 1
        li = []
        for slide in self.ppt.slides:
            li.append(["------slide " + str(i) + "---------","",""])
            for shape in slide.shapes:
                if shape.has_text_frame:
                    li.append([str(slide.slide_id),str(shape.shape_id), shape.text])
                if shape.has_table:
                    table = self._read_table(shape.table)
                    li.append([str(slide.slide_id),str(shape.shape_id),"table"])
                    for retu in table:
                        tex = ""
                        for moji in retu:
                            tex = tex + " , " + moji
                        li.append([slide.slide_id,shape.shape_id,tex])
            i += 1
        df_a = pd.DataFrame(li)
        df_a.columns=["slide_id","shape_id","text"]
        return df_a

    def _read_table(self,tbl):
        """tableのデータを読み込みます
        
        Arguments:
            tbl {shape} -- python-pptxのテーブルのshape
        """

        table_tex = []
        for r in enumerate(tbl.rows):
    #         print(r[0])
            temp = []
            for c in enumerate(tbl.columns):
    #             print(c[0])
                tex = tbl.cell(r[0],c[0]).text
                temp.append(tex)
            table_tex.append(temp)
        return table_tex


    def _find_shape(self,slide_id,shape_id):
        """slide_id（スライド固有）とshape_id(スライドの中のshape固有)からshapeを返します
        
        Arguments:
            slide_id {int} -- スライドの識別個別番号
            shape_id {int} -- shapeの識別番号（スライドの中では重複しない）
        
        Returns:
            shape -- python-pptxのshape
        """

        ppt = self.ppt
        for slide in ppt.slides:
            if slide.slide_id == slide_id:
                for shape in slide.shapes:
                    if shape.shape_id == shape_id:
                        tg = shape
                        return tg


if __name__ == '__main__':
    obj = pptx("test.pptx")
    df = pd.DataFrame(
        {'District':['Hampshire', 'Dorset', 'Wiltshire', 'Worcestershire'],
        'Population':[25000, 500000, 735298, 12653],
        'Ratio':[1.56, 7.34, 3.67, 8.23],
        'dddd':[15,65,25,65]})
    obj.add_table_from_df(df)
    obj.pptx_to_csv()
    obj.csv_to_pptx()

